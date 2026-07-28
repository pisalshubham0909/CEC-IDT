import http.server
import socketserver
import webbrowser
import threading
import time
import sys

PORT = 8000

class ThreadingHTTPServer(socketserver.ThreadingMixIn, http.server.HTTPServer):
    daemon_threads = True

def run_server():
    global PORT
    handler = http.server.SimpleHTTPRequestHandler
    
    class CORSRequestHandler(handler):
        def end_headers(self):
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Cache-Control', 'no-store, no-cache, must-revalidate, max-age=0')
            super().end_headers()

        def do_OPTIONS(self):
            self.send_response(200)
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Access-Control-Allow-Methods', 'POST, GET, OPTIONS')
            self.send_header('Access-Control-Allow-Headers', 'X-Password, X-Level, X-Mode, Content-Type, X-Session-ID')
            self.end_headers()

        def send_api_response(self, data, content_type="application/octet-stream"):
            self.send_response(200)
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Content-Type', content_type)
            self.send_header('Content-Length', str(len(data)))
            self.end_headers()
            self.wfile.write(data)

        def do_POST(self):
            if self.path.startswith("/api/"):
                self.handle_api()
            else:
                self.send_error(404)

        def handle_api(self):
            try:
                content_length = int(self.headers.get('Content-Length', 0))
                body = self.rfile.read(content_length) if content_length > 0 else b""
                
                import io, uuid, struct, tempfile, os
                try:
                    import pypdfium2 as pdfium
                    HAS_PDFIUM = True
                except ImportError:
                    HAS_PDFIUM = False
                from pypdf import PdfReader, PdfWriter

                global MERGE_SESSIONS
                if 'MERGE_SESSIONS' not in globals():
                    MERGE_SESSIONS = {}

                if self.path == "/api/merge_start":
                    session_id = str(uuid.uuid4())
                    MERGE_SESSIONS[session_id] = {
                        'master': pdfium.PdfDocument.new() if HAS_PDFIUM else PdfWriter(),
                        'type': 'pdfium' if HAS_PDFIUM else 'pypdf',
                        'created': time.time()
                    }
                    self.send_api_response(session_id.encode('utf-8'), "text/plain")

                elif self.path.startswith("/api/merge_chunk"):
                    session_id = self.headers.get('X-Session-ID', '')
                    if session_id in MERGE_SESSIONS:
                        sess = MERGE_SESSIONS[session_id]
                        offset = 0
                        body_len = len(body)
                        while offset < body_len:
                            if offset + 4 > body_len:
                                break
                            file_len = struct.unpack('>I', body[offset:offset+4])[0]
                            offset += 4
                            if offset + file_len > body_len:
                                break
                            file_bytes = body[offset:offset+file_len]
                            offset += file_len
                            
                            if sess['type'] == 'pdfium':
                                try:
                                    src_doc = pdfium.PdfDocument(file_bytes)
                                    sess['master'].import_pages(src_doc)
                                    src_doc.close()
                                except Exception as pdf_err:
                                    print(f"Warning during PDFium chunk merge: {pdf_err}")
                                    try:
                                        # Fallback to pypdf for problematic file
                                        reader = PdfReader(io.BytesIO(file_bytes), strict=False)
                                        temp_writer = PdfWriter()
                                        for p in reader.pages: temp_writer.add_page(p)
                                        t_buf = io.BytesIO()
                                        temp_writer.write(t_buf)
                                        temp_doc = pdfium.PdfDocument(t_buf.getvalue())
                                        sess['master'].import_pages(temp_doc)
                                        temp_doc.close()
                                    except Exception as e2:
                                        print(f"File skipped due to severe corruption: {e2}")
                            else:
                                try:
                                    sess['master'].append(io.BytesIO(file_bytes))
                                except Exception as pdf_err:
                                    print(f"Warning during pypdf chunk merge: {pdf_err}")

                        self.send_api_response(b"OK", "text/plain")
                    else:
                        self.send_error(400, "Invalid Session ID")

                elif self.path.startswith("/api/merge_finish"):
                    session_id = self.headers.get('X-Session-ID', '')
                    if session_id in MERGE_SESSIONS:
                        sess = MERGE_SESSIONS.pop(session_id)
                        out_buffer = io.BytesIO()
                        if sess['type'] == 'pdfium':
                            sess['master'].save(out_buffer)
                            sess['master'].close()
                        else:
                            sess['master'].write(out_buffer)
                        self.send_api_response(out_buffer.getvalue(), "application/pdf")
                    else:
                        self.send_error(400, "Invalid Session ID")

                elif self.path == "/api/encrypt":
                    try:
                        password = self.headers.get('X-Password', '')
                        out_buffer = io.BytesIO()
                        try:
                            import fitz
                            doc = fitz.open(stream=body, filetype="pdf")
                            perm = fitz.PDF_PERM_ACCESSIBILITY | fitz.PDF_PERM_PRINT | fitz.PDF_PERM_COPY | fitz.PDF_PERM_ANNOTATE | fitz.PDF_PERM_MODIFY
                            doc.save(out_buffer, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=password, owner_pw=password, permissions=perm)
                            doc.close()
                            res_bytes = out_buffer.getvalue()
                        except Exception as e_fitz:
                            print(f"PyMuPDF encrypt fallback note: {e_fitz}")
                            reader = PdfReader(io.BytesIO(body), strict=False)
                            writer = PdfWriter()
                            writer.append(reader)
                            writer.encrypt(user_password=password, owner_password=password)
                            writer.write(out_buffer)
                            res_bytes = out_buffer.getvalue()
                        
                        self.send_api_response(res_bytes, "application/pdf")
                    except Exception as enc_err:
                        print(f"Backend PDF encryption error: {enc_err}")
                        self.send_error(500, f"PDF encryption failed: {enc_err}")

                elif self.path == "/api/decrypt":
                    try:
                        password = self.headers.get('X-Password', '')
                        out_buffer = io.BytesIO()
                        try:
                            import fitz
                            doc = fitz.open(stream=body, filetype="pdf")
                            if doc.is_encrypted:
                                auth = doc.authenticate(password)
                                if not auth and password:
                                    doc.authenticate(password.encode('utf-8'))
                            doc.save(out_buffer, encryption=fitz.PDF_ENCRYPT_NONE)
                            doc.close()
                            res_bytes = out_buffer.getvalue()
                        except Exception as d_fitz:
                            print(f"PyMuPDF decrypt fallback note: {d_fitz}")
                            reader = PdfReader(io.BytesIO(body), strict=False)
                            if reader.is_encrypted:
                                reader.decrypt(password)
                            writer = PdfWriter()
                            writer.append(reader)
                            writer.write(out_buffer)
                            res_bytes = out_buffer.getvalue()
                            
                        self.send_api_response(res_bytes, "application/pdf")
                    except Exception as dec_err:
                        print(f"Backend PDF decryption error: {dec_err}")
                        self.send_error(500, f"PDF decryption failed: {dec_err}")
                    
                elif self.path == "/api/merge":
                    import struct
                    writer = PdfWriter()
                    offset = 0
                    body_len = len(body)
                    files_merged = 0
                    while offset < body_len:
                        if offset + 4 > body_len:
                            break
                        file_len = struct.unpack('>I', body[offset:offset+4])[0]
                        offset += 4
                        if offset + file_len > body_len:
                            break
                        file_bytes = body[offset:offset+file_len]
                        offset += file_len
                        try:
                            writer.append(io.BytesIO(file_bytes))
                            files_merged += 1
                        except Exception as pdf_err:
                            print(f"Warning during backend merge of file {files_merged+1}: {pdf_err}")

                    out_buffer = io.BytesIO()
                    writer.write(out_buffer)
                    self.send_api_response(out_buffer.getvalue(), "application/pdf")

                elif self.path == "/api/compress":
                    level = self.headers.get('X-Level', 'medium')
                    res_bytes = body
                    try:
                        try:
                            import fitz
                            HAS_FITZ = True
                        except ImportError:
                            HAS_FITZ = False

                        if HAS_FITZ:
                            doc = fitz.open(stream=body, filetype="pdf")
                            if level == 'high':
                                max_dim = 1000
                                quality = 35
                            elif level == 'low' or level == 'keep':
                                max_dim = 1600
                                quality = 65
                            else: # medium
                                max_dim = 1200
                                quality = 45

                            from PIL import Image
                            for page in doc:
                                for img_info in page.get_images(full=True):
                                    xref = img_info[0]
                                    try:
                                        base_image = doc.extract_image(xref)
                                        if base_image:
                                            image_bytes = base_image.get("image")
                                            if image_bytes:
                                                pil_img = Image.open(io.BytesIO(image_bytes))
                                                if pil_img.mode in ('RGBA', 'P', 'LA'):
                                                    pil_img = pil_img.convert('RGB')
                                                w, h = pil_img.size
                                                if max(w, h) > max_dim:
                                                    scale = max_dim / float(max(w, h))
                                                    new_w, new_h = max(1, int(w * scale)), max(1, int(h * scale))
                                                    pil_img = pil_img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                                                out_b = io.BytesIO()
                                                pil_img.save(out_b, format='JPEG', quality=quality, optimize=True)
                                                new_bytes = out_b.getvalue()
                                                try:
                                                    page.replace_image(xref, stream=new_bytes)
                                                except Exception:
                                                    doc.update_stream(xref, new_bytes)
                                    except Exception:
                                        pass
                            comp_bytes = doc.tobytes(garbage=4, deflate=True, clean=True, deflate_images=True, deflate_fonts=True)
                            if len(comp_bytes) < len(body):
                                res_bytes = comp_bytes
                            else:
                                res_bytes = comp_bytes
                        else:
                            reader = PdfReader(io.BytesIO(body), strict=False)
                            writer = PdfWriter()
                            writer.append(reader)

                            if level == 'high':
                                max_dim = 1200
                                quality = 45
                            elif level == 'low' or level == 'keep':
                                max_dim = 2400
                                quality = 80
                            else: # medium
                                max_dim = 1600
                                quality = 65

                            from PIL import Image
                            for page in writer.pages:
                                page.compress_content_streams()
                                try:
                                    for img_file in page.images:
                                        try:
                                            pil_img = img_file.image
                                            w, h = pil_img.size
                                            scale = 1.0
                                            if max(w, h) > max_dim:
                                                scale = max_dim / float(max(w, h))
                                                new_w, new_h = max(1, int(w * scale)), max(1, int(h * scale))
                                                pil_img = pil_img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                                            if pil_img.mode not in ('RGB', 'L'):
                                                pil_img = pil_img.convert('RGB')
                                            out_b = io.BytesIO()
                                            pil_img.save(out_b, format='JPEG', quality=quality, optimize=True)
                                            new_b = out_b.getvalue()
                                            if len(new_b) < len(img_file.data) or scale < 1.0:
                                                img_file.replace(pil_img, quality=quality)
                                        except Exception:
                                            pass
                                except Exception:
                                    pass

                            writer.compress_identical_objects(remove_duplicates=True, remove_unreferenced=True)
                            out_buffer = io.BytesIO()
                            writer.write(out_buffer)
                            comp_bytes = out_buffer.getvalue()
                            if len(comp_bytes) < len(body):
                                res_bytes = comp_bytes
                    except Exception as comp_err:
                        print(f"Backend compression fallback warning: {comp_err}")
                        res_bytes = body

                    self.send_api_response(res_bytes, "application/pdf")
                    
                elif self.path == "/api/convert_pdf_to_word":
                    try:
                        mode = self.headers.get('X-Mode', 'layout')
                        import fitz, docx
                        from docx.shared import Pt, Inches, RGBColor

                        doc_pdf = fitz.open('pdf', body)
                        doc_word = docx.Document()
                        
                        for page_idx in range(len(doc_pdf)):
                            p = doc_pdf[page_idx]
                            rect = p.rect
                            pw, ph = rect.width, rect.height
                            
                            section = doc_word.sections[0] if page_idx == 0 else doc_word.add_section()
                            section.page_width = Inches(pw / 72.0)
                            section.page_height = Inches(ph / 72.0)
                            section.top_margin = Inches(0)
                            section.bottom_margin = Inches(0)
                            section.left_margin = Inches(0)
                            section.right_margin = Inches(0)
                            
                            if mode == 'text':
                                blocks = p.get_text('dict')['blocks']
                                has_text = any(b.get('type') == 0 and len(b.get('lines', [])) > 0 for b in blocks)
                                if has_text:
                                    for b in blocks:
                                        if b.get('type') == 0:
                                            for line in b.get('lines', []):
                                                spans = line.get('spans', [])
                                                if not spans: continue
                                                p_para = doc_word.add_paragraph()
                                                p_para.paragraph_format.space_before = Pt(0)
                                                p_para.paragraph_format.space_after = Pt(2)
                                                p_para.paragraph_format.line_spacing = 1.15
                                                for sp in spans:
                                                    t_str = sp.get('text', '')
                                                    if not t_str: continue
                                                    run = p_para.add_run(t_str)
                                                    s_pt = sp.get('size', 10)
                                                    run.font.size = Pt(max(8, min(48, round(s_pt))))
                                                    flg = sp.get('flags', 0)
                                                    if flg & 2: run.font.bold = True
                                                    if flg & 1: run.font.italic = True
                                                    clr_int = sp.get('color', 0)
                                                    if clr_int != 0:
                                                        r_c = (clr_int >> 16) & 0xFF
                                                        g_c = (clr_int >> 8) & 0xFF
                                                        b_c = clr_int & 0xFF
                                                        run.font.color.rgb = RGBColor(r_c, g_c, b_c)
                                else:
                                    pix = p.get_pixmap(dpi=300)
                                    img_bytes = pix.tobytes('png')
                                    img_stream = io.BytesIO(img_bytes)
                                    p_para = doc_word.add_paragraph()
                                    p_para.paragraph_format.space_before = Inches(0)
                                    p_para.paragraph_format.space_after = Inches(0)
                                    p_para.add_run().add_picture(img_stream, width=Inches(pw / 72.0), height=Inches(ph / 72.0))
                            else:
                                # High-Fidelity Canvas Mode (100% Exact Visual Format Preservation)
                                pix = p.get_pixmap(dpi=300)
                                img_bytes = pix.tobytes('png')
                                img_stream = io.BytesIO(img_bytes)
                                p_para = doc_word.add_paragraph()
                                p_para.paragraph_format.space_before = Inches(0)
                                p_para.paragraph_format.space_after = Inches(0)
                                p_para.add_run().add_picture(img_stream, width=Inches(pw / 72.0), height=Inches(ph / 72.0))

                            if page_idx < len(doc_pdf) - 1:
                                doc_word.add_page_break()

                        out_b = io.BytesIO()
                        doc_word.save(out_b)
                        res_docx = out_b.getvalue()
                        self.send_api_response(res_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                    except Exception as conv_err:
                        print(f"Backend PDF to Word conversion error: {conv_err}")
                        self.send_error(500, f"PDF to Word conversion failed: {conv_err}")

                elif self.path == "/api/convert_pdf_to_pptx":
                    try:
                        mode = self.headers.get('X-Mode', 'layout')
                        import fitz
                        from pptx import Presentation
                        from pptx.util import Inches, Pt
                        from pptx.dml.color import RGBColor
                        
                        doc_pdf = fitz.open('pdf', body)
                        prs = Presentation()

                        for page_idx in range(len(doc_pdf)):
                            page = doc_pdf[page_idx]
                            rect = page.rect
                            pw, ph = rect.width, rect.height
                            
                            prs.slide_width = Inches(pw / 72.0)
                            prs.slide_height = Inches(ph / 72.0)
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            
                            if mode == 'text':
                                blocks = page.get_text('dict')['blocks']
                                has_text = any(b.get('type') == 0 and len(b.get('lines', [])) > 0 for b in blocks)
                                if has_text:
                                    for b in blocks:
                                        if b.get('type') == 0:
                                            bbox = b.get('bbox', (0,0,100,50))
                                            left = Inches(bbox[0] / 72.0)
                                            top = Inches(bbox[1] / 72.0)
                                            width = Inches(max(0.2, ((bbox[2] - bbox[0]) / 72.0) + 0.2))
                                            height = Inches(max(0.2, ((bbox[3] - bbox[1]) / 72.0) + 0.1))
                                            
                                            txBox = slide.shapes.add_textbox(left, top, width, height)
                                            tf = txBox.text_frame
                                            tf.word_wrap = True
                                            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
                                            
                                            para_idx = 0
                                            for line in b.get('lines', []):
                                                p_para = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
                                                para_idx += 1
                                                for span in line.get('spans', []):
                                                    txt = span.get('text', '')
                                                    if not txt: continue
                                                    run = p_para.add_run()
                                                    run.text = txt
                                                    size_pt = span.get('size', 10)
                                                    run.font.size = Pt(max(6, min(72, round(size_pt))))
                                                    flags = span.get('flags', 0)
                                                    if flags & 2: run.font.bold = True
                                                    if flags & 1: run.font.italic = True
                                                    color_int = span.get('color', 0)
                                                    if color_int != 0:
                                                        r = (color_int >> 16) & 0xFF
                                                        g = (color_int >> 8) & 0xFF
                                                        b_c = color_int & 0xFF
                                                        run.font.color.rgb = RGBColor(r, g, b_c)
                                else:
                                    pix = page.get_pixmap(dpi=300)
                                    img_bytes = pix.tobytes('png')
                                    img_stream = io.BytesIO(img_bytes)
                                    slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(pw / 72.0), Inches(ph / 72.0))
                            else:
                                # High-Fidelity Canvas Mode (100% Exact Visual Format Preservation)
                                pix = page.get_pixmap(dpi=300)
                                img_bytes = pix.tobytes('png')
                                img_stream = io.BytesIO(img_bytes)
                                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(pw / 72.0), Inches(ph / 72.0))

                        out_buf = io.BytesIO()
                        prs.save(out_buf)
                        self.send_api_response(out_buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                    except Exception as conv_err:
                        print(f"Backend PDF to PPTX conversion error: {conv_err}")
                        self.send_error(500, f"PDF to PPTX conversion failed: {conv_err}")

                elif self.path == "/api/convert_word_to_pdf":
                    try:
                        orientation = self.headers.get('X-PDF-Orientation', 'portrait')
                        theme_hex = self.headers.get('X-PDF-Theme', '#0284C7')
                        
                        import docx, html
                        from reportlab.lib.pagesizes import letter, landscape
                        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
                        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                        from reportlab.lib import colors

                        doc_in = docx.Document(io.BytesIO(body))
                        out_buf = io.BytesIO()
                        
                        pagesize_cfg = landscape(letter) if orientation == 'landscape' else letter
                        usable_w = (792 - 72) if orientation == 'landscape' else (612 - 72)
                        
                        pdf_doc = SimpleDocTemplate(out_buf, pagesize=pagesize_cfg, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
                        styles = getSampleStyleSheet()
                        normal_style = styles['Normal']
                        story = []

                        def _extract_images_from_run(r_element, doc_part):
                            imgs = []
                            try:
                                drawings = r_element._r.xpath('.//w:drawing')
                                for d in drawings:
                                    blips = d.xpath('.//a:blip/@r:embed')
                                    for rId in blips:
                                        if rId in doc_part.related_parts:
                                            img_data = doc_part.related_parts[rId].blob
                                            img_stream = io.BytesIO(img_data)
                                            imgs.append(RLImage(img_stream, width=min(usable_w, 400), height=200))
                            except Exception:
                                pass
                            return imgs

                        def _get_p_elements(p_item, doc_part):
                            runs_html = []
                            run_imgs = []
                            max_font_pt = 11
                            
                            for r in p_item.runs:
                                t = r.text
                                if t:
                                    t_html = html.escape(t)
                                    if r.bold: t_html = f'<b>{t_html}</b>'
                                    if r.italic: t_html = f'<i>{t_html}</i>'
                                    if r.font.color and r.font.color.rgb:
                                        hex_col = '#' + str(r.font.color.rgb)
                                        t_html = f'<font color="{hex_col}">{t_html}</font>'
                                    runs_html.append(t_html)
                                    if r.font.size and r.font.size.pt:
                                        max_font_pt = max(max_font_pt, int(r.font.size.pt))
                                
                                imgs = _extract_images_from_run(r, doc_part)
                                run_imgs.extend(imgs)
                            
                            elements = []
                            para_text = ''.join(runs_html) if runs_html else html.escape(p_item.text)
                            
                            if para_text.strip():
                                if p_item.style.name.startswith('Heading 1'): max_font_pt = max(max_font_pt, 20)
                                elif p_item.style.name.startswith('Heading 2'): max_font_pt = max(max_font_pt, 16)
                                elif p_item.style.name.startswith('Heading 3'): max_font_pt = max(max_font_pt, 13)
                                
                                align_code = 0
                                if p_item.alignment == docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER: align_code = 1
                                elif p_item.alignment == docx.enum.text.WD_ALIGN_PARAGRAPH.RIGHT: align_code = 2
                                elif p_item.alignment == docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY: align_code = 4

                                p_style = ParagraphStyle(
                                    'DocxParaStyle',
                                    parent=normal_style,
                                    fontSize=max_font_pt,
                                    leading=int(max_font_pt * 1.3),
                                    alignment=align_code,
                                    spaceAfter=6
                                )
                                elements.append(Paragraph(para_text, p_style))
                            
                            for img_item in run_imgs:
                                elements.append(img_item)
                                elements.append(Spacer(1, 6))
                                
                            return elements

                        def _get_tbl_element(tbl_item):
                            table_data = []
                            num_cols = max([len(row.cells) for row in tbl_item.rows]) if tbl_item.rows else 1
                            col_w = usable_w / float(num_cols)
                            
                            for row in tbl_item.rows:
                                row_cells = []
                                for cell in row.cells:
                                    cell_text = html.escape(cell.text.strip())
                                    cell_p = Paragraph(cell_text, ParagraphStyle('CellP', parent=normal_style, fontSize=10, leading=13))
                                    row_cells.append(cell_p)
                                table_data.append(row_cells)
                            if table_data:
                                t = Table(table_data, colWidths=[col_w] * num_cols)
                                t.setStyle(TableStyle([
                                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor(theme_hex)),
                                    ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                                    ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                                    ('TOPPADDING', (0,0), (-1,-1), 5),
                                    ('BOTTOMPADDING', (0,0), (-1,-1), 5),
                                ]))
                                return t
                            return None

                        for p_item in doc_in.paragraphs:
                            p_els = _get_p_elements(p_item, doc_in.part)
                            story.extend(p_els)
                            
                        for tbl_item in doc_in.tables:
                            t_el = _get_tbl_element(tbl_item)
                            if t_el:
                                story.append(t_el)
                                story.append(Spacer(1, 8))

                        if not story:
                            story.append(Paragraph('Document Content Empty', normal_style))

                        pdf_doc.build(story)
                        self.send_api_response(out_buf.getvalue(), "application/pdf")
                    except Exception as conv_err:
                        print(f"Backend Word to PDF conversion error: {conv_err}")
                        self.send_error(500, f"Word to PDF conversion failed: {conv_err}")
                    
                else:
                    self.send_error(404, "API endpoint not found")
            except Exception as e:
                self.send_response(500)
                self.send_header('Access-Control-Allow-Origin', '*')
                self.send_header('Content-Type', 'text/plain')
                self.end_headers()
                self.wfile.write(f"Error: {str(e)}".encode('utf-8'))

    # Find an open port starting from 8000
    for attempt in range(20):
        try:
            httpd = ThreadingHTTPServer(("", PORT), CORSRequestHandler)
            print(f"\n=======================================================")
            print(f"  CEC IDT Automation Tool Server is running!")
            print(f"  URL: http://localhost:{PORT}")
            print(f"  Press Ctrl+C to stop the server.")
            print(f"=======================================================\n")
            
            def open_browser():
                time.sleep(0.8)
                webbrowser.open(f"http://localhost:{PORT}")
            
            threading.Thread(target=open_browser, daemon=True).start()
            httpd.serve_forever()
            break
        except OSError:
            print(f"Port {PORT} is in use. Trying next port...")
            PORT += 1
    else:
        print("Could not find an available port to start the server. Exiting.")
        sys.exit(1)

if __name__ == "__main__":
    try:
        run_server()
    except KeyboardInterrupt:
        print("\nStopping server. Goodbye!")
        sys.exit(0)
